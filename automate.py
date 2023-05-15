from docx2pdf import convert
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path

try:
    docx_file = input("Enter a docx file: ")
    convert(docx_file, "output.pdf")
    pdf_file = 'output.pdf'
    print('Converted to PDF.')
except:
    print("File not found. Please enter the full path or put the file in the same directory.")

pages = convert_from_path('output.pdf', 500)

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

for count, page in enumerate(pages):
    page.save(f'out{count}.jpg', 'JPEG')

    slide = prs.slides.add_slide(blank_slide_layout)
    img_path = f'out{count}.jpg'
    pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(5))

prs.save('final.pptx')
print('Converted to PPT.')
